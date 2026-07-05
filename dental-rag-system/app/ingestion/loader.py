# ============================================================
# 文档加载器 — 把各种格式的文件解析成纯文本
# 支持：PDF、DOCX、TXT、Markdown
# 输出：LoadedDocument 对象（包含按页的文本内容）
# ============================================================
#导入依赖包

from __future__ import annotations #允许在类定义中引用自身

from pathlib import Path
from typing import Optional #可选类型，表示参数可以是某种类型或 None
from loguru import logger #比 print 更好的日志工具，能看到时间戳和模块名
from typing import List #列表类型，表示参数可以是某种类型的列表
import csv #用于处理 CSV 文件，提供了方便的接口来读取和写入 CSV 数据，这样就能更好地处理 CSV 文件了

#定义 Documentpage 类，表示加载后显示的文档
class DocumentPage:
    def __init__(self, content: str, page_number: int):
        self.content = content
        self.page_number = page_number #页码，从1开始

#定义 LoadedDocument 类，表示加载的文档
class LoadedDocument:
    def __init__(
        self,#先写注释
        pages: List[DocumentPage], #文档的页列表，每页是一个 DocumentPage 对象
        file_path: str, #文档的原始文件路径
        filename: str,  #文档的文件名
        metadata: Optional[dict] = None,
        ) ->None:#文档的元信息，可以是任何键值对，默认为 None
        self.pages = pages
        self.file_path = file_path
        self.filename = filename
        self.metadata = metadata or {} #如果没有提供元信息，就用一个空字典

    @property
    def full_text(self) -> str:
        """把所有页的文本拼成一段。"""
        return "\n".join(p.content for p in self.pages)
    
    @property
    def char_count(self) -> int:
        """计算总的字符数"""
        return sum(len(p.content) for p in self.pages)

    @property
    def page_count(self) -> int:
        """计算总的页数"""
        return len(self.pages)
    
#定义 Loader 类，表示文档加载器
# (LOADERS dict moved after function definitions) #键是文件扩展名，值是对应的加载函数，这样的写法更好，更容易维护和扩展

def load_document(file_path: str) -> LoadedDocument:
    path = Path(file_path)
    ext = path.suffix.lower()
    
    # 直接从注册表查找，效率更高，扩展性更好
    loader = LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"Unsupported file type: {ext}")
        
    return loader(file_path, path.name)




def _load_pdf(file_path: str, filename: str) -> LoadedDocument: #这里使用Path对象作为注释，后续代码格式更规范
    """加载 PDF 文件。"""
    #需要做try except 来捕获可能的错误，比如pdf环境不存在、格式错误等
    try:
        import pypdf
    except ImportError:
        raise ImportError("pypdf library is required to load PDF files. Please install it with 'pip install pypdf'.")
    # pages = [] #type: List[DocumentPage]

    pages: List[DocumentPage] = [] #更规范的写法，明确类型
    with open(file_path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for i, page in enumerate(reader.pages):
            text = (page.extract_text() or "").strip() #如果提取文本失败，就用空字符串,strip()去掉首尾空白符,这点很重要
            if text: #如果文本不为空才添加到页面列表中，这样可以避免空页，但是就无法获取真实的页面的页数
                pages.append(DocumentPage(content=text, page_number=i + 1))
    return LoadedDocument(pages=pages, file_path=str(file_path), filename=filename)#最终的返回值是一个 LoadedDocument 对象，包含了所有页面的文本内容和相关信息



def _load_docx(file_path: str, filename: str) -> LoadedDocument:
    """加载 DOCX 文件。"""
    try:
        import docx
    except ImportError:
        raise ImportError("python-docx library is required to load DOCX files. Please install it with 'pip install python-docx'.")
    
    doc = docx.Document(file_path)
    all_text:List[str] = []
    #用于处理docx中的普通的文本段落，docx库会把每个段落当成一个对象，我们需要遍历这些段落来提取文本内容
    for para in doc.paragraphs:
        if para.text.strip(): #只保留非空段落
            all_text.append(para.text.strip())

    #用于处理docx中的表格，表格中的文本也很重要，我们需要遍历所有的表格和单元格来提取文本内容
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip(): #只保留非空单元格
                    all_text.append(cell.text.strip())

    #把所有的文本内容拼接成一个字符串，作为文档的唯一页面，这样就能保留所有的文本信息了
    full_text = "\n".join(all_text)if all_text else "" #如果没有任何文本，就用一个空字符串，避免返回 None
    pages = [DocumentPage(content=full_text, page_number=1)] #把所有文本作为一个页面

    return LoadedDocument(pages=pages, file_path=str(file_path), filename=filename)
        

def _load_text(file_path: str, filename: str) -> LoadedDocument:
    """加载 TXT 文件。"""
    #txt文件的加载相对简单，我们只需要直接读取文件内容即可，注意要指定编码和错误处理方式，以避免因为编码问题导致的加载失败
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    pages = [DocumentPage(content=text, page_number=1)] 
    return LoadedDocument(pages=pages, file_path=str(file_path), filename=filename)

def _load_markdown(file_path: str, filename: str) -> LoadedDocument:
    """加载 Markdown 文件。"""
    
    try:
        from bs4 import BeautifulSoup
        import markdown
    except ImportError:
        raise ImportError("markdown library is required to load Markdown files. Please install it with 'pip install markdown'.")
    
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    
    #把 Markdown 转换成 HTML，然后再提取文本内容，这样就能保留一些 Markdown 的结构信息了
    html = markdown.markdown(text)
    soup = BeautifulSoup(html, "html.parser")
    clean_text = soup.get_text(separator="\n").strip() #去掉 HTML 标签，得到干净的文本字符串
    pages = [DocumentPage(content=clean_text, page_number=1)]
    return LoadedDocument(pages=pages, file_path=str(file_path), filename=filename)

def _load_csv(file_path: str, filename: str) -> LoadedDocument:
    """加载 CSV 文件。"""

    #csv文件的加载稍微复杂一些，我们需要先读取文件内容，然后解析每一行，把每一行转换成一个文本行，列之间用逗号分隔，这样就能保留表格的结构信息了

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)
    
        if not rows: #如果没有任何行，就返回一个空的 LoadedDocument 对象
            return LoadedDocument(pages=[], file_path=str(file_path), filename=filename) #如果没有任何行，就返回一个空的 LoadedDocument 对象

        #把表格的每一行都转换成一个文本行，列之间用逗号分隔，这样就能保留表格的结构信息了
        text_rows = [", ".join(row) for row in rows if any(cell.strip() for cell in row)] #只保留非空行，避免空行干扰
        full_text = "\n".join(text_rows) if text_rows else "" #如果没有任何文本，就用一个空字符串
        pages = [DocumentPage(content=full_text, page_number=1)] #把所有文本作为一个页面

    return LoadedDocument(pages=pages, file_path=str(file_path), filename=filename)

# def _load_xlsx(file_path: str, filename: str) -> LoadedDocument:
#     """加载 XLSX 文件。"""
#     try:
#         import pandas as pd

#     except ImportError:
#         raise ImportError("pandas library is required to load XLSX files. Please install it with 'pip install pandas'.")
#     try:
#         df = pd.read_excel(file_path)#使用 pandas 读取 Excel 文件，得到一个 DataFrame 对象，DataFrame 是一种表格数据结构，可以方便地进行数据处理和分析
#         if df.empty: #如果 DataFrame 为空，就返回一个空的 LoadedDocument 对象
#             raise ValueError("Excel file is empty.")
#     except FileNotFoundError:
#         raise ValueError(f"Failed to read Excel file: {file_path} not found.")
#     except Exception as e:
#         raise ValueError(f"Failed to read Excel file: {e}")
    
#     #把 DataFrame 的每一行都转换成一个文本行，列之间用逗号分隔，这样就能保留表格的结构信息了
#     text_rows = [", ".join(map(str, row)) for row in df.values]
#     full_text = "\n".join(text_rows) if text_rows else ""
#     pages = [DocumentPage(content=full_text, page_number=1)]
#     return LoadedDocument(pages=pages, file_path=str(file_path), filename=filename)


def _load_xlsx(file_path: str, filename: str) -> LoadedDocument:
    """加载 XLSX 文件。"""
    try:
        import pandas as pd
    except ImportError:
        raise ImportError("pandas library is required to load XLSX files")
    
    try:
        df = pd.read_excel(file_path)
        if df.empty:
            raise ValueError("Excel file is empty or has no data")
    except Exception as e:
        raise ValueError(f"Failed to read Excel file: {str(e)}")
    
    # 基本的文本清理
    def clean_cell(val):
        return "" if pd.isna(val) else str(val).strip()#如果单元格为空，就返回空字符串，否则返回单元格的字符串表示，并去掉前后的空白字符
    
    headers = ", ".join(clean_cell(col) for col in df.columns)#把列名也作为第一行文本，这样就能保留表头信息了
    text_rows = [headers]
    
    for _, row in df.iterrows():
        cleaned_row = [clean_cell(val) for val in row]
        text_rows.append(", ".join(cleaned_row))
    
    full_text = "\n".join(text_rows)
    pages = [DocumentPage(content=full_text, page_number=1)]
    
    return LoadedDocument(
        file_path=str(file_path),
        filename=filename,
        pages=pages
    )



def _load_pptx(file_path: str, filename: str) -> LoadedDocument:
    """加载 PPTX 文件。"""
    try:
        import pptx
    except ImportError:
        raise ImportError("python-pptx library is required to load PPTX files. Please install it with 'pip install python-pptx'.")


    doc = pptx.Presentation(file_path)
    all_text: List[str] = []
    for slide in doc.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    all_text.append(text)
    full_text = "\n".join(all_text) if all_text else ""
    pages = [DocumentPage(content=full_text, page_number=1)]
    return LoadedDocument(pages=pages, file_path=str(file_path), filename=filename)

def _load_html(file_path: str, filename: str) -> LoadedDocument:
    """加载 HTML 文件。"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("beautifulsoup4 library is required to load HTML files. Please install it with 'pip install beautifulsoup4'.")
    
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f, "html.parser")
    
    #提取所有文本内容，去掉多余的空白符，这样就能得到一个干净的文本字符串了
    text = soup.get_text(separator="\n").strip()
    pages = []
    for i, page_text in enumerate(text.split("\n\n")): #把文本按双换行分成多页，这样就能保留一些结构信息了
        page_text = page_text.strip()
        if page_text: #只保留非空页
            pages.append(DocumentPage(content=page_text, page_number=i + 1))
    return LoadedDocument(pages=pages, file_path=str(file_path), filename=filename)

        

# 全局加载器注册表
LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".txt": _load_text,
    ".md": _load_markdown,
    ".csv": _load_csv,
    ".xlsx": _load_xlsx,
    ".pptx": _load_pptx,
    ".html": _load_html,
}





