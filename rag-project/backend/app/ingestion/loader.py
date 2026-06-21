# ============================================================
# 文档加载器 — 把各种格式的文件解析成纯文本
# 支持：PDF、DOCX、TXT、Markdown
# 输出：LoadedDocument 对象（包含按页的文本内容）
# ============================================================

from __future__ import annotations

from pathlib import Path
from typing import Optional

from loguru import logger  # 比 print 更好的日志工具，能看到时间戳和模块名


class DocumentPage:
    """代表文档中的一页。"""
    def __init__(self, text: str, page_number: int) -> None:
        self.text = text             # 这一页的纯文本内容
        self.page_number = page_number  # 页码（从1开始）


class LoadedDocument:
    """解析后的完整文档。"""
    def __init__(
        self,
        file_path: str,               # 源文件路径
        filename: str,                # 文件名
        pages: list[DocumentPage],    # 所有页的列表
        metadata: Optional[dict] = None,  # 额外的元信息
    ) -> None:
        self.file_path = file_path
        self.filename = filename
        self.pages = pages
        self.metadata = metadata or {}

    @property
    def full_text(self) -> str:
        """把所有页的文本拼成一段。"""
        return "\n".join(p.text for p in self.pages)

    @property
    def char_count(self) -> int:
        """总字符数。"""
        return sum(len(p.text) for p in self.pages)

    @property
    def page_count(self) -> int:
        """总页数。"""
        return len(self.pages)


def load_document(file_path: str) -> LoadedDocument:
    """根据文件后缀名自动选择加载器。"""
    path = Path(file_path) #把字符串路径转换成 Path 对象，方便操作
    ext = path.suffix.lower()       # 获取扩展名，如 .pdf
    logger.info("Loading: {} ({})", path.name, ext)  # 打印日志（延时计算、节约性能）

    if ext == ".pdf":
        return _load_pdf(file_path, path.name)
    elif ext == ".docx":
        return _load_docx(file_path, path.name)
    elif ext in (".txt", ".md"):
        return _load_text(file_path, path.name)
    elif ext == ".csv":
        return _load_csv(file_path, path.name)
    elif ext == ".xlsx":
        return _load_xlsx(file_path, path.name)
    elif ext in (".jpg", ".jpeg", ".png", ".bmp", ".gif"): 
        return _load_image(file_path, path.name)
    elif ext == ".pptx":
        return _load_pptx(file_path, path.name)
    elif ext == ".html":
        return _load_html(file_path, path.name)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ── PDF 解析 ────────────────────────────────────────────

def _load_pdf(file_path: str, filename: str) -> LoadedDocument:
    """使用 pypdf 库解析 PDF 文件，按页提取文本。"""
    try:
        import pypdf
    except ImportError:
        raise ImportError("pip install pypdf")  # 没装库就提示安装

    pages: list[DocumentPage] = []
    with open(file_path, "rb") as f:  # rb = 二进制读取
        reader = pypdf.PdfReader(f)    # 创建 PDF 读取器
        for i, page in enumerate(reader.pages):
            # 提取文本，如果提取不到就用空字符串
            text = (page.extract_text() or "").strip()
            if text:  # 只保留非空页
                pages.append(DocumentPage(text=text, page_number=i + 1))
    return LoadedDocument(file_path, filename, pages)


# ── DOCX 解析 ───────────────────────────────────────────

def _load_docx(file_path: str, filename: str) -> LoadedDocument:
    """使用 python-docx 库解析 Word 文档。"""
    try:
        import docx
    except ImportError:
        raise ImportError("pip install python-docx")

    doc = docx.Document(file_path)
    # 提取所有段落，去掉空行
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs) if paragraphs else ""
    # Word 没有"页"的概念，所以所有文本放在一页里
    pages = [DocumentPage(text=text, page_number=1)]
    return LoadedDocument(file_path, filename, pages)


# ── TXT / Markdown 解析 ────────────────────────────────

def _load_text(file_path: str, filename: str) -> LoadedDocument:
    """普通文本文件，直接全部读取。"""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        # errors="replace" 遇到编码错误也不崩溃，替换成 �
        text = f.read()
    pages = [DocumentPage(text=text, page_number=1)]
    return LoadedDocument(file_path, filename, pages)



# ── CSV 解析 ───────────────────────────────────────────

def _load_csv(file_path: str, filename: str) -> LoadedDocument:
    """使用内置 csv 库解析 CSV 文件，将表格转为文本行。"""
    import csv
    
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)
        
        if not rows:
            return LoadedDocument(file_path, filename, [])
            
        # 把每一行用逗号拼起来，行与行之间换行
        # 例如: ["姓名", "年龄"] -> "姓名, 年龄"
        text_lines = [", ".join(row) for row in rows]
        text = "\n".join(text_lines)
        
    pages = [DocumentPage(text=text, page_number=1)]
    return LoadedDocument(file_path, filename, pages)

# Excel (.xlsx) 同理，可以用 openpyxl 库，逻辑一样，只是读取方式变成 cell.value


# ── Excel (.xlsx) 解析 ───────────────────────────────────────────
def _load_xlsx(file_path: str, filename: str) -> LoadedDocument:
    """使用 openpyxl 库解析 Excel 文件，将表格转为文本行。"""
    try:
        import openpyxl
    except ImportError:
        raise ImportError("pip install openpyxl")

    # read_only=True 极大节省内存，特别是处理大文件时
    wb = openpyxl.load_workbook(file_path, read_only=True)
    ws = wb.active  # 默认只读取第一个工作表
    
    text_lines = []
    try:
        for row in ws.iter_rows(values_only=True):
            # 1. 将单元格转为字符串，None 转为空字符串
            cells = [str(cell) if cell is not None else "" for cell in row]
            
            # 2. 过滤掉全空的行（比如表格里啥都没写的行）
            if any(cell.strip() for cell in cells):
                line = ", ".join(cells)
                text_lines.append(line)
    finally:
        # 3. 关键修改：read_only 模式必须手动关闭，否则会内存泄漏！
        wb.close()
        
    text = "\n".join(text_lines)
    pages = [DocumentPage(text=text, page_number=1)]
    return LoadedDocument(file_path, filename, pages)





# ── 图片解析 ───────────────────────────────────────────

def _load_image(file_path: str, filename: str) -> LoadedDocument:
    """使用 pytesseract 库识别图片中的文字。"""
    try:
        import pytesseract
    except ImportError:
        raise ImportError("请安装 OCR 依赖: pip install pytesseract Pillow")

    from PIL import Image

    try:
        with Image.open(file_path) as img:
            # 关键修改 1：指定识别语言。
            # 'chi_sim+eng' 表示同时支持简体中文和英文。如果只识别英文，请改为 'eng'
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            
    except pytesseract.TesseractNotFoundError:
        # 关键修改 2：捕获 Tesseract 引擎未安装的错误，给出友好提示
        raise RuntimeError(
            "Tesseract OCR 引擎未安装或未加入系统环境变量。\n"
            "请根据你的操作系统安装 Tesseract：\n"
            "  Windows: 下载安装包并添加到 PATH\n"
            "  Mac: brew install tesseract\n"
            "  Linux: sudo apt install tesseract-ocr"
        )
    except Exception as e:
        # 关键修改 3：捕获图片损坏等未知错误，防止整个程序崩溃
        logger.warning("Failed to OCR image {}: {}", filename, e)
        text = ""  # 识别失败则置为空文本

    pages = [DocumentPage(text=text, page_number=1)]
    return LoadedDocument(file_path, filename, pages)




# ── PPT 解析 ───────────────────────────────────────────
def _load_pptx(file_path: str, filename: str) -> LoadedDocument:
    """使用 python-pptx 解析 PPT，按幻灯片分页。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("pip install python-pptx")

    pages: list[DocumentPage] = []
    prs = Presentation(file_path)
    
    for i, slide in enumerate(prs.slides):
        texts = []
        # 一个 slide 里有多个 shape（文本框、图片等），只提取有文本的
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        
        if texts: # 如果这页有文字
            text = "\n".join(texts)
            pages.append(DocumentPage(text=text, page_number=i + 1))
            
    return LoadedDocument(file_path, filename, pages)


# ── HTML 解析 ───────────────────────────────────────────

def _load_html(file_path: str, filename: str) -> LoadedDocument:
    """使用 BeautifulSoup 解析 HTML，只保留纯文本。"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("pip install beautifulsoup4")

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f, "html.parser")
        
        # 可选：去掉 script 和 css 标签内容
        for script in soup(["script", "style"]):
            script.extract()
            
    text = soup.get_text(separator="\n", strip=True) # 获取纯文本并自动去首尾空格
    pages = [DocumentPage(text=text, page_number=1)]
    return LoadedDocument(file_path, filename, pages)




